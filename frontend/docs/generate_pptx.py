from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Brand colors
DARK_TEAL   = RGBColor(0x22, 0x38, 0x43)   # #223843
LIGHT_BG    = RGBColor(0xEF, 0xF1, 0xF3)   # #eff1f3
MAUVE       = RGBColor(0xDB, 0xD3, 0xD8)   # #dbd3d8
SAND        = RGBColor(0xD8, 0xB4, 0xA0)   # #d8b4a0
TERRACOTTA  = RGBColor(0xD7, 0x7A, 0x61)   # #d77a61
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

# ─────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1, DARK_TEAL)

add_rect(s1, 0, 5.8, 13.33, 1.7, TERRACOTTA)

add_text(s1, "Travel Management App",
         0.7, 1.6, 11, 1.2,
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text(s1, "Missing Features & Backend Blockers",
         0.7, 2.85, 11, 0.8,
         font_size=22, bold=False, color=SAND, align=PP_ALIGN.LEFT)

add_text(s1, "Prepared for Engineering Handoff  •  March 2026",
         0.7, 6.1, 11, 0.5,
         font_size=13, color=WHITE, italic=True, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────
# Slide 2 — Current Status
# ─────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2, LIGHT_BG)

add_rect(s2, 0, 0, 13.33, 1.15, DARK_TEAL)
add_text(s2, "Current Implementation Status",
         0.5, 0.18, 12, 0.8,
         font_size=26, bold=True, color=WHITE)

steps = [
    ("Step 1", "API Client + Auth JWT + RBAC Guards", TERRACOTTA),
    ("Step 2", "Traveler Flow  (catalog → detail → subscription → tracking)", TERRACOTTA),
    ("Step 3", "Manager Flow  (CRUD travels + subscribers)", TERRACOTTA),
    ("Step 4", "Admin Flow  (users + notifications + payments)", TERRACOTTA),
    ("Step 5", "Advanced Search  (Elasticsearch filters)", TERRACOTTA),
    ("Step 6", "Backlog  (analytics, reviews, recommendations)", MAUVE),
]

for i, (label, desc, color) in enumerate(steps):
    y = 1.4 + i * 0.88
    add_rect(s2, 0.5, y, 1.5, 0.62, color)
    add_text(s2, label, 0.5, y + 0.1, 1.5, 0.45,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    status = "DONE" if color == TERRACOTTA else "TODO"
    add_text(s2, desc + f"  [{status}]",
             2.3, y + 0.1, 10.5, 0.5,
             font_size=14, bold=(status == "TODO"), color=BLACK if status == "DONE" else TERRACOTTA)

# ─────────────────────────────────────────────
# Slide 3 — Buildable Now
# ─────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3, LIGHT_BG)

add_rect(s3, 0, 0, 13.33, 1.15, DARK_TEAL)
add_text(s3, "Can Be Built Now  (Endpoints Exist, No UI Yet)",
         0.5, 0.18, 12, 0.8,
         font_size=26, bold=True, color=WHITE)

add_rect(s3, 0.5, 1.35, 0.18, 4.8, TERRACOTTA)

items = [
    ("User Profile Page",
     ["Edit profile  →  PUT /api/v1/users/me",
      "Change password  →  POST /api/v1/users/me/change-password",
      "Endpoint ready — only the UI page is missing"]),
    ("Manager — Payments per Travel",
     ["View all payments for a travel  →  GET /api/v1/payments/travel/{travelId}",
      "Listed in manager spec but no page was built",
      "Endpoint ready — only the UI page is missing"]),
]

y = 1.45
for title, bullets in items:
    add_text(s3, title, 0.9, y, 11.5, 0.5,
             font_size=18, bold=True, color=DARK_TEAL)
    y += 0.48
    for b in bullets:
        add_text(s3, f"  •  {b}", 1.0, y, 11.2, 0.38,
                 font_size=13, color=BLACK)
        y += 0.37
    y += 0.2

# ─────────────────────────────────────────────
# Slide 4 — Blocked items
# ─────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4, LIGHT_BG)

add_rect(s4, 0, 0, 13.33, 1.15, DARK_TEAL)
add_text(s4, "Blocked — Backend Endpoints Required",
         0.5, 0.18, 12, 0.8,
         font_size=26, bold=True, color=WHITE)

blocked = [
    ("Admin Analytics Dashboard",
     "No aggregation endpoints. Needs: /api/v1/admin/dashboard or equivalent."),
    ("Manager Analytics Dashboard",
     "No aggregate endpoint. Needs revenue, traveler count, performance breakdown."),
    ("Reviews & Ratings",
     "No review service at all. Needs a full review/report microservice + REST controllers."),
    ("Recommendations (Neo4j)",
     "rec-service exists but has ZERO REST controllers. Needs /api/v1/recommendations/me."),
]

for i, (title, detail) in enumerate(blocked):
    y = 1.4 + i * 1.35
    add_rect(s4, 0.5, y, 12.3, 1.15, WHITE, MAUVE)
    add_rect(s4, 0.5, y, 0.22, 1.15, TERRACOTTA)
    add_text(s4, title, 0.9, y + 0.08, 11.2, 0.42,
             font_size=16, bold=True, color=DARK_TEAL)
    add_text(s4, detail, 0.9, y + 0.55, 11.2, 0.5,
             font_size=12, color=BLACK)

# ─────────────────────────────────────────────
# Slide 5 — Summary table
# ─────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5, LIGHT_BG)

add_rect(s5, 0, 0, 13.33, 1.15, DARK_TEAL)
add_text(s5, "Summary Table",
         0.5, 0.18, 12, 0.8,
         font_size=26, bold=True, color=WHITE)

headers = ["Feature", "Backend Endpoint", "Frontend UI", "Status"]
col_x   = [0.4, 4.2, 7.6, 10.5]
col_w   = [3.7, 3.3, 2.8, 2.7]

add_rect(s5, 0.4, 1.25, 12.5, 0.52, DARK_TEAL)
for j, h in enumerate(headers):
    add_text(s5, h, col_x[j] + 0.08, 1.3, col_w[j], 0.42,
             font_size=13, bold=True, color=WHITE)

rows = [
    ("Profile edit + change password", "Exists", "Not built", "Build now"),
    ("Manager payments view",          "Exists", "Not built", "Build now"),
    ("Admin analytics dashboard",      "Missing","Not built", "Backend needed"),
    ("Manager analytics dashboard",    "Missing","Not built", "Backend needed"),
    ("Reviews & ratings",              "Missing","Not built", "Backend needed"),
    ("Recommendations (Neo4j)",        "Missing","Not built", "Backend needed"),
]

for i, row in enumerate(rows):
    y = 1.82 + i * 0.72
    row_bg = WHITE if i % 2 == 0 else LIGHT_BG
    add_rect(s5, 0.4, y, 12.5, 0.65, row_bg)
    status_color = TERRACOTTA if "needed" in row[3] else RGBColor(0x2E, 0x7D, 0x32)
    for j, cell in enumerate(row):
        c = status_color if j == 3 else BLACK
        bold = j == 3
        add_text(s5, cell, col_x[j] + 0.08, y + 0.1, col_w[j], 0.48,
                 font_size=12, bold=bold, color=c)

# ─────────────────────────────────────────────
# Slide 6 — Action Items
# ─────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6, DARK_TEAL)

add_rect(s6, 0, 0, 13.33, 1.15, TERRACOTTA)
add_text(s6, "Backend Action Items",
         0.5, 0.18, 12, 0.8,
         font_size=26, bold=True, color=WHITE)

actions = [
    ("1", "Expose Admin analytics endpoint",
     "Aggregate revenue, top managers, monthly activity — e.g. GET /api/v1/admin/dashboard"),
    ("2", "Expose Manager analytics endpoint",
     "Aggregate revenue + traveler count per manager — e.g. GET /api/v1/manager/dashboard"),
    ("3", "Implement Review service",
     "New microservice with controllers for POST /reviews, GET /reviews/travel/{id}, report endpoint"),
    ("4", "Implement Recommendation REST controllers",
     "rec-service already exists in architecture — add controllers for GET /api/v1/recommendations/me"),
]

for i, (num, title, detail) in enumerate(actions):
    y = 1.4 + i * 1.35
    add_rect(s6, 0.5, y, 0.6, 1.1, TERRACOTTA)
    add_text(s6, num, 0.5, y + 0.2, 0.6, 0.65,
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s6, 1.2, y, 11.6, 1.1, RGBColor(0x2E, 0x4A, 0x56))
    add_text(s6, title, 1.35, y + 0.08, 11.2, 0.42,
             font_size=16, bold=True, color=SAND)
    add_text(s6, detail, 1.35, y + 0.55, 11.2, 0.48,
             font_size=12, color=LIGHT_BG)

# ─────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────
out = "/home/cheikh/Documents/frontend/docs/Travel_App_Missing_Features.pptx"
prs.save(out)
print(f"Saved: {out}")
